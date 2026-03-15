from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime, timedelta

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(88, 86, 214)  # Purple
SECONDARY_COLOR = RGBColor(255, 107, 107)  # Red
ACCENT_COLOR = RGBColor(74, 144, 226)  # Blue
DARK_BG = RGBColor(30, 30, 40)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 240, 245)

HEADER_DARK = RGBColor(20, 20, 28)
MUTED_TEXT = RGBColor(120, 120, 130)

SOFT_BG_1 = RGBColor(246, 247, 252)
SOFT_BG_2 = RGBColor(238, 241, 255)
INK = RGBColor(18, 18, 28)

def _set_slide_bg(slide, rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def add_bg_accents(slide):
    blob1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.8), Inches(1.2), Inches(4.0), Inches(4.0))
    blob1.fill.solid(); blob1.fill.fore_color.rgb = SOFT_BG_2
    blob1.line.fill.background()
    blob2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), Inches(-1.2), Inches(4.2), Inches(4.2))
    blob2.fill.solid(); blob2.fill.fore_color.rgb = RGBColor(255, 240, 245)
    blob2.line.fill.background()
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.65), Inches(10), Inches(0.85))
    band.fill.solid(); band.fill.fore_color.rgb = SOFT_BG_1
    band.line.fill.background()

def add_header(slide, title="AI Component Generator", subtitle=""):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.65))
    header.fill.solid()
    header.fill.fore_color.rgb = HEADER_DARK
    header.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(6.8), Inches(0.5))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        s = slide.shapes.add_textbox(Inches(7.2), Inches(0.14), Inches(2.6), Inches(0.5))
        sf = s.text_frame
        sp = sf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(10)
        sp.font.color.rgb = LIGHT_GRAY
        sp.alignment = PP_ALIGN.RIGHT

def add_footer(slide, left_text="", right_text=""):
    return

def add_section_chip(slide, text, x, y, color):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.2), Inches(0.42))
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    tf = chip.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_feature_card(slide, x, y, w, h, title, body, accent_rgb):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(250, 250, 252)
    card.line.color.rgb = RGBColor(225, 225, 235)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_rgb
    bar.line.fill.background()

    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.25), Inches(y + 0.25), Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = accent_rgb
    dot.line.fill.background()

    t = slide.shapes.add_textbox(Inches(x + 0.75), Inches(y + 0.18), Inches(w - 1.0), Inches(0.5))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(35, 35, 45)

    b = slide.shapes.add_textbox(Inches(x + 0.75), Inches(y + 0.72), Inches(w - 1.0), Inches(h - 0.85))
    bf = b.text_frame
    bp = bf.paragraphs[0]
    bp.text = body
    bp.font.size = Pt(12)
    bp.font.color.rgb = RGBColor(70, 70, 80)
    bp.line_spacing = 1.15

def add_key_features_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    add_bg_accents(slide)
    add_header(slide, subtitle="Highlights")
    add_section_chip(slide, "KEY FEATURES", 0.5, 0.9, PRIMARY_COLOR)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "What makes GenUI useful"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = INK

    add_feature_card(
        slide, 0.5, 2.2, 4.6, 1.55,
        "AI Code Generation",
        "Turn a short UI description into clean component code using backend AI generation.",
        PRIMARY_COLOR,
    )
    add_feature_card(
        slide, 5.0, 2.2, 4.5, 1.55,
        "Live Preview + Editor",
        "Generated code opens inside Monaco Editor with instant preview for rapid iteration.",
        ACCENT_COLOR,
    )
    add_feature_card(
        slide, 0.5, 3.95, 4.6, 1.55,
        "Protected Access",
        "Protected routes and session token ensure controlled access to features.",
        RGBColor(76, 175, 80),
    )
    add_feature_card(
        slide, 5.0, 3.95, 4.5, 1.55,
        "Supabase History",
        "Every generation is stored in Supabase so users can revisit recent prompts and outputs.",
        SECONDARY_COLOR,
    )
    # No footer / watermark

def add_architecture_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    add_bg_accents(slide)
    add_header(slide, subtitle="System")
    add_section_chip(slide, "ARCHITECTURE", 0.5, 0.9, ACCENT_COLOR)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "High-level architecture"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 25, 35)

    client = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.4), Inches(2.7), Inches(1.25))
    client.fill.solid(); client.fill.fore_color.rgb = RGBColor(250, 250, 252)
    client.line.color.rgb = RGBColor(230, 230, 240)
    client.text_frame.text = "Frontend\nReact + Vite"
    client.text_frame.paragraphs[0].font.bold = True

    api = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.9), Inches(2.4), Inches(2.7), Inches(1.25))
    api.fill.solid(); api.fill.fore_color.rgb = RGBColor(250, 250, 252)
    api.line.color.rgb = RGBColor(230, 230, 240)
    api.text_frame.text = "Backend API\nExpress"
    api.text_frame.paragraphs[0].font.bold = True

    gemini = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(2.1), Inches(2.2), Inches(1.1))
    gemini.fill.solid(); gemini.fill.fore_color.rgb = RGBColor(76, 175, 80)
    gemini.line.fill.background()
    gemini.text_frame.text = "Gemini\nAPI"
    gemini.text_frame.paragraphs[0].font.bold = True
    gemini.text_frame.paragraphs[0].font.color.rgb = WHITE

    supa = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(3.5), Inches(2.2), Inches(1.1))
    supa.fill.solid(); supa.fill.fore_color.rgb = PRIMARY_COLOR
    supa.line.fill.background()
    supa.text_frame.text = "Supabase\nDB"
    supa.text_frame.paragraphs[0].font.bold = True
    supa.text_frame.paragraphs[0].font.color.rgb = WHITE

    a1 = slide.shapes.add_connector(1, Inches(3.4), Inches(3.0), Inches(3.9), Inches(3.0))
    a1.line.color.rgb = RGBColor(90, 90, 100)
    a2 = slide.shapes.add_connector(1, Inches(6.6), Inches(2.9), Inches(7.1), Inches(2.7))
    a2.line.color.rgb = RGBColor(90, 90, 100)
    a3 = slide.shapes.add_connector(1, Inches(6.6), Inches(3.15), Inches(7.1), Inches(4.05))
    a3.line.color.rgb = RGBColor(90, 90, 100)

    note = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(6.0), Inches(2.4))
    nf = note.text_frame
    np = nf.paragraphs[0]
    np.text = "Data flow:\n• Frontend calls /generate (Bearer token)\n• Backend calls Gemini for code\n• Backend stores output in Supabase history\n• Frontend fetches /history to display recent generations"
    np.font.size = Pt(13)
    np.font.color.rgb = RGBColor(70, 70, 80)
    np.line_spacing = 1.2

    # No footer / watermark

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide, PRIMARY_COLOR)

    accent = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.0), Inches(-1.2), Inches(5.0), Inches(5.0))
    accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor(255, 255, 255)
    accent.fill.fore_color.brightness = -0.85
    accent.line.fill.background()
    accent2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), Inches(3.2), Inches(4.5), Inches(4.5))
    accent2.fill.solid(); accent2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    accent2.fill.fore_color.brightness = -0.88
    accent2.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Footer
    # No footer / watermark

def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide, WHITE)
    add_bg_accents(slide)
    add_header(slide, subtitle="Documentation")
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = INK
    
    # Add colored line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.75), Inches(9), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_COLOR
    
    # Add content
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.1), Inches(9), Inches(4.9))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(230, 230, 240)
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.35), Inches(8.4), Inches(4.6))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(8)
        p.space_after = Pt(8)

    # No footer / watermark

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    add_bg_accents(slide)
    add_header(slide, subtitle="Architecture")
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = INK
    
    # Line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.75), Inches(9), Inches(0))
    line.line.color.rgb = PRIMARY_COLOR
    line.line.width = Pt(3)
    
    # Left column
    left_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.1), Inches(4.6), Inches(4.9))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = RGBColor(250, 250, 252)
    left_panel.line.color.rgb = RGBColor(230, 230, 240)
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.35), Inches(4.0), Inches(4.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_content):
        if i > 0:
            p = left_frame.add_paragraph()
        else:
            p = left_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(6)
    
    # Right column
    right_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(2.1), Inches(4.5), Inches(4.9))
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = RGBColor(250, 250, 252)
    right_panel.line.color.rgb = RGBColor(230, 230, 240)
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(2.35), Inches(3.9), Inches(4.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_content):
        if i > 0:
            p = right_frame.add_paragraph()
        else:
            p = right_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(6)

    # No footer / watermark

# Slide 1: Title Slide
add_title_slide(prs, "AI Component Generator", "Intelligent Code Generation Platform")

# New Slide: Key Features (visual)
add_key_features_slide(prs)

# Slide 2: Introduction
add_content_slide(prs, "Introduction & Objectives", [
    "• Objective: Automate frontend component generation using AI",
    "• Leverage Google Gemini API for intelligent code synthesis",
    "• Support multiple frameworks: HTML, CSS, Tailwind, Bootstrap",
    "• Enable developers to rapidly prototype UI components",
    "• Reduce development time and improve code consistency",
    "• Provide real-time preview and code editing capabilities"
])

# New Slide: Architecture Overview (visual)
add_architecture_overview_slide(prs)

# Slide 3: System Requirements Specification (SRS)
add_content_slide(prs, "System Requirement Specification (SRS)", [
    "Functional Requirements:",
    "  • User input framework selection (HTML/CSS/Tailwind/Bootstrap)",
    "  • Component description input via textarea",
    "  • API integration with Google Generative AI",
    "  • Code editor with syntax highlighting",
    "  • Live preview of generated components",
    "  • Copy-to-clipboard functionality",
    "  • Persist generated code history (Supabase) per authenticated session",
    "  • Error handling and fallback mechanisms"
])

# Slide 4: Process Logic
add_content_slide(prs, "Process Logic", [
    "1. User selects framework from dropdown menu",
    "2. User describes desired component in textarea",
    "3. Frontend sends request to backend API (/generate)",
    "4. Backend constructs prompt for Gemini model",
    "5. Gemini API generates component code",
    "6. Backend returns generated code to frontend",
    "7. Frontend displays code in Monaco editor",
    "8. User can preview, copy, or export the code"
])

# Slide 5: Gantt Chart
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Project Timeline (3 Months)"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = PRIMARY_COLOR

line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
line.line.color.rgb = PRIMARY_COLOR
line.line.width = Pt(3)

# Gantt chart data
gantt_data = [
    ("Requirements & Design", 1, 3, ACCENT_COLOR),
    ("Backend Development", 2, 5, SECONDARY_COLOR),
    ("Frontend Development", 3, 6, RGBColor(76, 175, 80)),
    ("API Integration", 4, 7, RGBColor(255, 152, 0)),
    ("Testing & QA", 6, 8, RGBColor(156, 39, 176)),
    ("Deployment", 9, 9, PRIMARY_COLOR),
]

y_pos = 2.0
for task, start, end, color in gantt_data:
    # Task name
    task_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.4))
    task_frame = task_box.text_frame
    p = task_frame.paragraphs[0]
    p.text = task
    p.font.size = Pt(13)
    p.font.bold = True
    
    # Bar
    bar_width = (end - start) * 0.6
    bar = slide.shapes.add_shape(1, Inches(2.8 + start * 0.6), Inches(y_pos + 0.05), Inches(bar_width), Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.color.rgb = color
    
    y_pos += 0.6

# Legend
legend_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
legend_frame = legend_box.text_frame
p = legend_frame.paragraphs[0]
p.text = "Timeline: Week 1-9  |  3 Months Total Duration"
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = RGBColor(100, 100, 100)

# Slide 6: Data Dictionary
add_content_slide(prs, "Data Dictionary", [
    "Framework: string - selected UI framework (HTML/CSS/Tailwind/Bootstrap)",
    "Prompt: string - user's component description",
    "Code: string - generated component code output",
    "Loading: boolean - indicates API request in progress",
    "OutputScreen: boolean - toggles code editor/preview visibility",
    "Tab: integer - current tab selection (1=Code, 2=Preview)",
    "Response: object - API response containing generated code",
    "Error: string - error message if generation fails",
    "History: array - list of previously generated components loaded from Supabase",
    "HistoryItem: object - { id, framework, prompt, code, created_at }",
    "AuthToken: string - bearer token used to associate history with a user/session"
])

# Slide 7: DFD (Data Flow Diagram) - updated with proper shapes
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Data Flow Diagram (DFD)"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = PRIMARY_COLOR

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(0.06))
line.fill.solid()
line.fill.fore_color.rgb = PRIMARY_COLOR

# External Entities (rectangles)
user_ext = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.0), Inches(1.6), Inches(0.8))
user_ext.fill.solid()
user_ext.fill.fore_color.rgb = ACCENT_COLOR
user_ext.line.color.rgb = ACCENT_COLOR
user_ext.text_frame.text = "User"
user_ext.text_frame.paragraphs[0].font.color.rgb = WHITE
user_ext.text_frame.paragraphs[0].font.bold = True

templates_ext = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(2.2), Inches(0.9))
templates_ext.fill.solid()
templates_ext.fill.fore_color.rgb = RGBColor(245,245,245)
templates_ext.line.color.rgb = RGBColor(120,120,120)
templates_ext.text_frame.text = "Templates Store"
templates_ext.text_frame.paragraphs[0].font.bold = True

# Processes (circles/ovals)
proc_validate = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.0), Inches(2.0), Inches(1.4), Inches(1.0))
proc_validate.fill.solid()
proc_validate.fill.fore_color.rgb = PRIMARY_COLOR
proc_validate.line.color.rgb = PRIMARY_COLOR
proc_validate.text_frame.text = "Validate\nPrompt"
proc_validate.text_frame.paragraphs[0].font.color.rgb = WHITE
proc_validate.text_frame.paragraphs[0].font.bold = True

proc_generate = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.0), Inches(2.0), Inches(1.6), Inches(1.0))
proc_generate.fill.solid()
proc_generate.fill.fore_color.rgb = SECONDARY_COLOR
proc_generate.line.color.rgb = SECONDARY_COLOR
proc_generate.text_frame.text = "Generate\nComponent"
proc_generate.text_frame.paragraphs[0].font.color.rgb = WHITE
proc_generate.text_frame.paragraphs[0].font.bold = True

# Data Stores (open-ended rectangles simulated by overlay)
ds_templates = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(4.5), Inches(2.4), Inches(1.0))
ds_templates.fill.background()
ds_templates.line.color.rgb = RGBColor(80,80,80)
ds_templates.text_frame.text = "Framework\nTemplates"
ds_templates.text_frame.paragraphs[0].font.bold = True

# Make right side of data store 'open' by overlaying a white rectangle to hide the border
open_mask = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.35), Inches(4.5), Inches(0.9), Inches(1.0))
open_mask.fill.solid()
open_mask.fill.fore_color.rgb = WHITE
open_mask.line.fill.background()

ds_cache = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(4.5), Inches(2.6), Inches(1.0))
ds_cache.fill.background()
ds_cache.line.color.rgb = RGBColor(80,80,80)
ds_cache.text_frame.text = "Cache &\nHistory"
ds_cache.text_frame.paragraphs[0].font.bold = True

open_mask2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(4.5), Inches(0.6), Inches(1.0))
open_mask2.fill.solid()
open_mask2.fill.fore_color.rgb = WHITE
open_mask2.line.fill.background()

# External API (rectangle)
api_ext = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(2.0), Inches(1.6), Inches(0.8))
api_ext.fill.solid()
api_ext.fill.fore_color.rgb = RGBColor(76,175,80)
api_ext.line.color.rgb = RGBColor(76,175,80)
api_ext.text_frame.text = "Gemini API"
api_ext.text_frame.paragraphs[0].font.color.rgb = WHITE
api_ext.text_frame.paragraphs[0].font.bold = True

# Arrows (connectors)
conn1 = slide.shapes.add_connector(1, Inches(2.1), Inches(2.4), Inches(3.0), Inches(2.5))
conn1.line.color.rgb = RGBColor(80,80,80)
conn2 = slide.shapes.add_connector(1, Inches(4.4), Inches(2.5), Inches(5.0), Inches(2.5))
conn2.line.color.rgb = RGBColor(80,80,80)
conn3 = slide.shapes.add_connector(1, Inches(6.6), Inches(2.5), Inches(7.8), Inches(2.5))
conn3.line.color.rgb = RGBColor(80,80,80)

conn4 = slide.shapes.add_connector(1, Inches(5.8), Inches(3.0), Inches(6.6), Inches(4.0))
conn4.line.color.rgb = RGBColor(80,80,80)
conn5 = slide.shapes.add_connector(1, Inches(4.2), Inches(4.0), Inches(3.0), Inches(4.6))
conn5.line.color.rgb = RGBColor(80,80,80)

# Labels for arrows
lbl1 = slide.shapes.add_textbox(Inches(2.1), Inches(2.0), Inches(1.2), Inches(0.3))
lbl1.text_frame.text = "Framework\nPrompt"
lbl2 = slide.shapes.add_textbox(Inches(4.4), Inches(2.0), Inches(1.2), Inches(0.3))
lbl2.text_frame.text = "Request"
lbl3 = slide.shapes.add_textbox(Inches(6.6), Inches(2.0), Inches(1.2), Inches(0.3))
lbl3.text_frame.text = "Prompt"

# Footer
footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.4))
fframe = footer.text_frame
f = fframe.paragraphs[0]
f.text = "DFD: External Entities (rectangles), Processes (circles), Data Stores (open rectangles)"
f.font.size = Pt(10)
f.font.color.rgb = LIGHT_GRAY
f.alignment = PP_ALIGN.LEFT

# Slide 8: ER Diagram
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Entity-Relationship Diagram (ER)"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = PRIMARY_COLOR

line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
line.line.color.rgb = PRIMARY_COLOR
line.line.width = Pt(3)

# User Entity
user_entity = slide.shapes.add_shape(1, Inches(0.8), Inches(2), Inches(1.5), Inches(1.2))
user_entity.fill.solid()
user_entity.fill.fore_color.rgb = ACCENT_COLOR
text_frame = user_entity.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "USER\n—\nuser_id\nname\nemail"
p.font.size = Pt(11)
p.font.color.rgb = WHITE
p.font.bold = True

# Generation Entity
gen_entity = slide.shapes.add_shape(1, Inches(3.5), Inches(2), Inches(1.8), Inches(1.2))
gen_entity.fill.solid()
gen_entity.fill.fore_color.rgb = PRIMARY_COLOR
text_frame = gen_entity.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "GENERATION\n—\ngen_id\nframework\nprompt"
p.font.size = Pt(11)
p.font.color.rgb = WHITE
p.font.bold = True

# Code Entity
code_entity = slide.shapes.add_shape(1, Inches(6.2), Inches(2), Inches(1.5), Inches(1.2))
code_entity.fill.solid()
code_entity.fill.fore_color.rgb = SECONDARY_COLOR
text_frame = code_entity.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "CODE\n—\ncode_id\ngenerated\ntimestamp"
p.font.size = Pt(11)
p.font.color.rgb = WHITE
p.font.bold = True

# API Entity
api_entity = slide.shapes.add_shape(1, Inches(8), Inches(2), Inches(1.5), Inches(1.2))
api_entity.fill.solid()
api_entity.fill.fore_color.rgb = RGBColor(76, 175, 80)
text_frame = api_entity.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "API_LOG\n—\napi_id\nstatus\nresponse"
p.font.size = Pt(11)
p.font.color.rgb = WHITE
p.font.bold = True

# Relationships
rel1 = slide.shapes.add_connector(1, Inches(2.3), Inches(2.6), Inches(3.5), Inches(2.6))
rel1.line.color.rgb = RGBColor(100, 100, 100)
rel1_label = slide.shapes.add_textbox(Inches(2.5), Inches(2.2), Inches(1), Inches(0.3))
rel1_frame = rel1_label.text_frame
p = rel1_frame.paragraphs[0]
p.text = "creates"
p.font.size = Pt(10)

rel2 = slide.shapes.add_connector(1, Inches(5.3), Inches(2.6), Inches(6.2), Inches(2.6))
rel2.line.color.rgb = RGBColor(100, 100, 100)
rel2_label = slide.shapes.add_textbox(Inches(5.5), Inches(2.2), Inches(1), Inches(0.3))
rel2_frame = rel2_label.text_frame
p = rel2_frame.paragraphs[0]
p.text = "stores"
p.font.size = Pt(10)

rel3 = slide.shapes.add_connector(1, Inches(7.7), Inches(2.6), Inches(8), Inches(2.6))
rel3.line.color.rgb = RGBColor(100, 100, 100)
rel3_label = slide.shapes.add_textbox(Inches(7.6), Inches(2.2), Inches(1.2), Inches(0.3))
rel3_frame = rel3_label.text_frame
p = rel3_frame.paragraphs[0]
p.text = "logs"
p.font.size = Pt(10)

# Relationships info
rel_info = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(8.4), Inches(3))
rel_frame = rel_info.text_frame
rel_frame.word_wrap = True
p = rel_frame.paragraphs[0]
p.text = "Relationships:\n• USER creates multiple GENERATIONs (1:N)\n• GENERATION results in CODE (1:1)\n• CODE logs into API_LOG (1:1)\n• USER can view their generation history"
p.font.size = Pt(13)
p.line_spacing = 1.3

# Slide 9: Interface
add_content_slide(prs, "User Interface Overview", [
    "Navigation Bar:",
    "  • Project logo and branding",
    "Framework Selection:",
    "  • Dropdown with 5 framework options",
    "Component Description:",
    "  • Textarea for detailed component requirements",
    "Generate Button:",
    "  • Triggers API request with loading spinner",
    "Code Output:",
    "  • Monaco Editor with syntax highlighting",
    "Output Controls:",
    "  • Copy, Export, Preview, Refresh buttons"
])

# Slide 10: Expected Report Generation
add_content_slide(prs, "Expected Report Generation", [
    "Generated Component Report includes:",
    "  • Framework-specific code structure",
    "  • HTML markup and semantic tags",
    "  • CSS styling and responsive design",
    "  • JavaScript functionality (if applicable)",
    "  • Commentary and best practices",
    "Export Formats:",
    "  • Raw HTML/CSS files",
    "  • Downloadable code snippets",
    "  • Shareable component links"
])

# Slide 11: References & Bibliography
add_content_slide(prs, "References & Bibliography", [
    "Technologies & Frameworks:",
    "• Google Generative AI (Gemini API) - model-serving & inference",
    "• React.js 19.2.0 - frontend UI framework",
    "• Express.js 5.2.1 - backend API server",
    "• Vite 7.2.4 - build tool and dev server",
    "• Tailwind CSS 3.4.19 - utility-first styling",
    "• Monaco Editor - code editor component",
    "Documentation & Resources:",
    "• https://react.dev/link/react-devtools",
    "• https://google.ai/gemini"
])

# Slide 12: Future Scope
add_content_slide(prs, "Future Scope & Enhancements", [
    "Advanced Features:",
    "  • Support for additional frameworks (Vue, Angular, Svelte)",
    "  • Multi-language code generation (TypeScript, JSX, Python)",
    "  • Component templates and design system integration",
    "  • Version control and component history tracking",
    "  • Real-time collaboration for teams",
    "  • Custom AI model fine-tuning",
    "  • Mobile app for iOS/Android",
    "  • Cloud deployment with CDN integration"
])

# Save presentation (new file)
output_path = r"D:\ai-com-gen\AI_Component_Generator_PPT_Pro.pptx"
prs.save(output_path)
print(f"✅ PowerPoint presentation created successfully!")
print(f"📍 Location: {output_path}")
